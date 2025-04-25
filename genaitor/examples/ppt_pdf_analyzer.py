import asyncio
import PyPDF2
import pptx
from dotenv import load_dotenv

from core import Orchestrator, Flow, ExecutionMode
from presets.agents import document_agent, question_agent, search_agent, response_agent

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a PDF file."""
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        return text
    except Exception as e:
        return str(e)

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extracts text from a PPT file."""
    try:
        presentation = pptx.Presentation(ppt_path)
        text = "\n".join([
            shape.text for slide in presentation.slides
            for shape in slide.shapes if hasattr(shape, "text")
        ])
        return text
    except Exception as e:
        return str(e)

async def main():
    print("\nInitializing Document QA System...")

    orchestrator = Orchestrator(
        agents={
            "document_agent": document_agent,
            "question_agent": question_agent,
            "search_agent": search_agent,
            "response_agent": response_agent
        },
        flows={
            "document_qa_flow": Flow(
                agents=["document_agent", "question_agent", "search_agent", "response_agent"],
                context_pass=[True, True, True, True]
            )
        },
        mode=ExecutionMode.SEQUENTIAL
    )

    # Test document path
    document_path = r"examples\files\2308.08468v1.pdf"
    question = "What are the key takeaways from this document?"

    print(f"\nExtracting content from: {document_path}")

    if document_path.endswith(".pdf"):
        document_content = extract_text_from_pdf(document_path)
    elif document_path.endswith(".pptx"):
        document_content = extract_text_from_ppt(document_path)
    else:
        print("\nUnsupported document format.")
        return

    # Execute document Q&A flow
    print("\nStarting document question answering...")

    try:
        result = await orchestrator.process_request(
            {"document_content": document_content, "user_question": question},
            flow_name='document_qa_flow'
        )

        if result["success"]:
            print("\nGenerated Answer:")
            print(result['content']['response_agent'].content)
        else:
            print(f"\nError: {result['error']}")

    except Exception as e:
        print(f"\nError: {str(e)}")

if __name__ == "__main__":
    asyncio.run(main())
