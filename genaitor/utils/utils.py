from PIL import Image
import pytesseract
import os
import json
import pdfplumber
from pptx import Presentation
from docx import Document
import pandas as pd
import speech_recognition as sr
from moviepy.editor import AudioFileClip

from genaitor.config import config

def extract_text_from_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        return ' '.join(page.extract_text() for page in pdf.pages)

def extract_text_from_ppt(file_path):
    ppt = Presentation(file_path)
    return '\n'.join([shape.text for slide in ppt.slides for shape in slide.shapes if shape.has_text_frame])

def extract_text_from_doc(file_path):
    doc = Document(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_json(file_path):
    with open(file_path, 'r') as f:
        return json.dumps(json.load(f), indent=2)

def transcribe_audio_file(file_path):
    recognizer = sr.Recognizer()
    audio_format = os.path.splitext(file_path)[1].lower()
    if audio_format == '.mp4':
        audio = AudioFileClip(file_path)
        temp_audio_path = "temp_audio.wav"
        audio.write_audiofile(temp_audio_path)
        file_path = temp_audio_path

    with sr.AudioFile(file_path) as source:
        audio_data = recognizer.record(source)
    return recognizer.recognize_google(audio_data)

def read_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)

def image_to_text(image_path):
    # TODO: Add preprocess with FURIO before OCR.
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)