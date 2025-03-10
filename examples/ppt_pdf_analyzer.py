import os
import sys
import json
import asyncio
import PyPDF2
import pptx
from dotenv import load_dotenv
from typing import Dict, Any

# Add project path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from src.core import (
    Agent, Task, Orchestrator, Flow,
    ExecutionMode, AgentRole, TaskResult
)
from src.llm import GeminiProvider, GeminiConfig

# Load environment variables
load_dotenv('.env')

class AnalysisTask(Task):
    def __init__(self, description: str, goal: str, output_format: str, llm_provider):
        super().__init__(description, goal, output_format)
        self.llm = llm_provider

    def execute(self, input_data: Dict[str, Any]) -> TaskResult:
        prompt = f"""
        Task: {self.description}
        Goal: {self.goal}

        {json.dumps(input_data, indent=4)}

        Provide the response in the following format:
        {self.output_format}
        """

        try:
            response = self.llm.generate(prompt)
            return TaskResult(
                success=True,
                content=response,
                metadata={"task_type": self.description}
            )
        except Exception as e:
            return TaskResult(
                success=False,
                content=None,
                error=str(e)
            )

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a PDF file."""
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        return text
    except Exception as e:
        return str(e)

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extracts text from a PPT file."""
    try:
        presentation = pptx.Presentation(ppt_path)
        text = "\n".join([
            shape.text for slide in presentation.slides
            for shape in slide.shapes if hasattr(shape, "text")
        ])
        return text
    except Exception as e:
        return str(e)

async def main():
    print("\nInitializing Document QA System...")

    test_keys = [os.getenv('API_KEY')]

    # Configure Gemini LLM
    gemini_config = GeminiConfig(
        api_keys=test_keys,
        temperature=0.7,
        verbose=False,
        max_tokens=5000
    )
    provider = GeminiProvider(gemini_config)

    # Define tasks
    document_analysis = AnalysisTask(
        description="Document Analysis",
        goal="Extract relevant content from PDF or PPT documents",
        output_format="JSON format with extracted content",
        llm_provider=provider
    )

    question_analysis = AnalysisTask(
        description="Question Analysis",
        goal="Interpret the user's question and prepare it for content search",
        output_format="JSON format with the question and intended answer structure",
        llm_provider=provider
    )

    answer_search = AnalysisTask(
        description="Answer Search",
        goal="Search for the answer within the extracted content from the document",
        output_format="JSON format with the search result",
        llm_provider=provider
    )

    response_generation = AnalysisTask(
        description="Response Generation",
        goal="Generate a final response based on the found answer and provide it in a clear format",
        output_format="Final answer in text format",
        llm_provider=provider
    )

    # Create agents
    document_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[document_analysis],
        llm_provider=provider
    )

    question_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[question_analysis],
        llm_provider=provider
    )

    search_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[answer_search],
        llm_provider=provider
    )

    response_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[response_generation],
        llm_provider=provider
    )

    # Orchestrate tasks
    orchestrator = Orchestrator(
        agents={
            "document_agent": document_agent,
            "question_agent": question_agent,
            "search_agent": search_agent,
            "response_agent": response_agent
        },
        flows={
            "document_qa_flow": Flow(
                agents=["document_agent", "question_agent", "search_agent", "response_agent"],
                context_pass=[True, True, True, True]
            )
        },
        mode=ExecutionMode.SEQUENTIAL
    )

    # Test document path
    document_path = r"examples\files\2308.08468v1.pdf"
    question = "What are the key takeaways from this document?"

    print(f"\nExtracting content from: {document_path}")

    if document_path.endswith(".pdf"):
        document_content = extract_text_from_pdf(document_path)
    elif document_path.endswith(".pptx"):
        document_content = extract_text_from_ppt(document_path)
    else:
        print("\nUnsupported document format.")
        return

    # Execute document Q&A flow
    print("\nStarting document question answering...")

    try:
        result = await orchestrator.process_request(
            {"document_content": document_content, "user_question": question},
            flow_name='document_qa_flow'
        )

        if result["success"]:
            print("\nGenerated Answer:")
            print(result['content']['response_agent'].content)
        else:
            print(f"\nError: {result['error']}")

    except Exception as e:
        print(f"\nError: {str(e)}")

if __name__ == "__main__":
    asyncio.run(main())
