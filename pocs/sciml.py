from bs4 import BeautifulSoup
import requests
import streamlit as st

import tempfile
import asyncio
import sys
import os

import google.generativeai as genai
import pandas as pd
import json
import tifffile
from PIL import Image
import PyPDF2
from pptx import Presentation
import cv2
import scipy
import ezdxf
import numpy as np
import h5py
from netCDF4 import Dataset
from astropy.io import fits
import avro.datafile
import avro.schema
import trimesh
import open3d as o3d

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from src.genaitor.core import (
    Agent, Task, Orchestrator, TaskResult,
    ExecutionMode, AgentRole, Flow
)
from src.genaitor.llm import GeminiProvider, GeminiConfig

frameworks = {
        "DiffEqFlux.jl":"https://docs.sciml.ai/DiffEqFlux/stable/",
        "TensorFlow Probability":"https://www.tensorflow.org/probability?hl=pt-br",
        "PyMC": "https://www.pymc.io/welcome.html",
        "Scipy": "https://scipy.org/",
        "PySPH": "https://pythonhosted.org/PySPH/overview.html",
        "GPy": "https://gpy.readthedocs.io/en/deploy/",
        "Diffeqpy":"https://pypi.org/project/diffeqpy/",
        "Physics-Informed Neural Operators":"https://arxiv.org/pdf/2111.03794",
        "Variational Autoencoders": "https://www.geeksforgeeks.org/variational-autoencoders/",
        "Generative Adversarial Networks":"https://arxiv.org/pdf/1611.01673",
        "Reinforcement Learninig":"https://arxiv.org/pdf/2503.18892",
        "Symbolic Regression": "https://github.com/MilesCranmer/PySR",
        "Data Assimilation": "https://arxiv.org/pdf/2406.00390",
        "Uncertainty Quantification (UQ)":"https://en.wikipedia.org/wiki/Uncertainty_quantification",
        "DeepXDE":"https://deepxde.readthedocs.io/en/latest/",
        "Surrogate Models":"https://smt.readthedocs.io/en/latest/index.html",
        "Neural Operators":"https://neuraloperator.github.io/dev/user_guide/neural_operators.html",
        "Neuromancer": "https://github.com/pnnl/neuromancer",
        "Simulai": "https://github.com/IBM/simulai/",
        "TFC": "https://github.com/leakec/tfc",
        "FEM/FVM/FDM":"https://www.machinedesign.com/additive-3d-printing/fea-and-simulation/article/21832072/whats-the-difference-between-fem-fdm-and-fvm",
        "Physics Informed Extreme Learning Machine": "https://arxiv.org/pdf/1907.03507",
        "Bayesian Physics Informed Neural Networks": "https://arxiv.org/pdf/2205.08304",
        "Quantum Physics Informed Neural Networks": "https://arxiv.org/pdf/2503.12244",
        "Graph PINN's": "https://arxiv.org/pdf/2211.05520",
        "DeepXDE": "https://github.com/lululxvi/deepxde",
        "NVidia Modulus": "https://github.com/NVIDIA/physicsnemo",
        "OpenFOAM": "https://doc.openfoam.com/2312/quickstart/",
        "PhyGeoNet": "https://github.com/Jianxun-Wang/phygeonet",
        "SciML Julia": "https://docs.sciml.ai/Overview/stable/",
        "Runge-Kutta PINNs": "https://runge-kutta-pinn-cmi-cmi-public-aaca66b00a12c20ef6731763c12508c.pages.desy.de/",
        "Google DeepMind GraphCast": "https://github.com/google-deepmind/graphcast",
        "Classical PINNs": "https://github.com/maziarraissi/PINNs",
        "XTFC": "https://arxiv.org/pdf/2005.10632",
        "ROM":"https://pymor.org/",
        "Molecular Dynamics":"https://en.wikipedia.org/wiki/Molecular_dynamics",
        }

framework_names = ', '.join(list(frameworks.keys()))+'...'
frameworks_list = list(frameworks.keys())
frameworks_list.sort()

class SolutionChooserTask(Task):
    def __init__(self, description: str, goal: str, output_format: str, llm_provider):
        super().__init__(description, goal, output_format)
        self.llm = llm_provider

    def execute(self, input_data: str) -> TaskResult:
        prompt = f"""
        Task: {self.description}
        Goal: {self.goal}
        
        Based on different solutions options, choose the best solution for answer the user request

        {input_data}
        
        {self.output_format}
        """
        
        try:
            response = self.llm.generate(prompt)
            return TaskResult(
                success=True,
                content=response,
                metadata={"task_type": "training_code"}
            )
        except Exception as e:
            return TaskResult(
                success=False,
                content=None,
                error=str(e)
            )

class ContextExpanderTask(Task):
    def __init__(self, description: str, goal: str, output_format: str, llm_provider):
        super().__init__(description, goal, output_format)
        self.llm = llm_provider

    def execute(self, input_data: str) -> TaskResult:
        prompt = f"""
        Task: {self.description}
        Goal: {self.goal}
        
        Based on a simple context of a problem, expand it the most.

        Context: {input_data}
        {self.output_format}
        """
        
        try:
            response = self.llm.generate(prompt)
            return TaskResult(
                success=True,
                content=response,
                metadata={"task_type": "training_code"}
            )
        except Exception as e:
            return TaskResult(
                success=False,
                content=None,
                error=str(e)
            )

class CodeSolutionTask(Task):
    def __init__(self, description: str, goal: str, output_format: str, llm_provider):
        super().__init__(description, goal, output_format)
        self.llm = llm_provider

    def execute(self, input_data: str) -> TaskResult:
        prompt = f"""
        Task: {self.description}
        Goal: {self.goal}
        
        Based on a the context of a problem and on a framework to solve it, apply the specified framework to solve the problem.
        
        Ensure the code includes:
        - Best practices for model training and scientific machine learning
        - Error handling and logging mechanisms
        - Comments explaining each section of the code
        - Scientific explanation for the solution provided

        Also, create and fill all the functions and code, don't let anything open for the user.
        
        {input_data}
        {self.output_format}
        """
        
        try:
            response = self.llm.generate(prompt)
            return TaskResult(
                success=True,
                content=response,
                metadata={"task_type": "training_code"}
            )
        except Exception as e:
            return TaskResult(
                success=False,
                content=None,
                error=str(e)
            )

async def context_expansion(organized_context, provider):
    expander_context_task = ContextExpanderTask(
        description="Expand the context of a scientific problem.",
        goal="Given a simple description of a scientific problem, generate a more robust context of it.",
        output_format="The robust context of problem, keeping same main context.",
        llm_provider=provider
    )
    
    expander_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[expander_context_task],
        llm_provider=provider
    )
    
    orchestrator = Orchestrator(
        agents={"expander_agent": expander_agent},
        flows={
            "expander_context_flow": Flow(agents=["expander_agent"], context_pass=[True])
        },
        mode=ExecutionMode.SEQUENTIAL
    )
    
    expanded_context = await orchestrator.process_request(organized_context, flow_name="expander_context_flow")
    return expanded_context['content']['expander_agent'].content

async def choice_strategy(expanded_context, provider):
    
    strategies = [
        "Understanding Problem",
        "Exploratory Data Analysis for Scientific Machine Learning",
        "Scikit-Learn and Numpy Modelling",
        "Numerical Solvers: Monte Carlo, FFT, LU Decomposition, Conjugate Gradient, etc",
        "Deep Learning Solvers: Tensorflow, PyTorch, JAX, etc",
        f"Other Frameworks: {framework_names}"
    ]

    input_data = f"Context: \n{expanded_context}\nStrategies: \n{strategies}"

    strategy_choice_task = SolutionChooserTask(
        description="Choose a path to solve a scientific problem, based on the problem context.",
        goal=(
            "Given the context and different approaches for solving the problem, choose the one that is the best approach for solving the problem."
            "Notice that you only has access to this information, nothing else or external."    
        ),
        output_format=(
            "The Framework choosed to solve the problem."
            "If you think the answer should be one of the other frameworks, please answer with: A different framework: {name_of_framework}"
            "The answer should end with: The choosed strategy was: "
        ),
        llm_provider=provider
    )
    
    strategy_choice_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[strategy_choice_task],
        llm_provider=provider
    )
    
    orchestrator = Orchestrator(
        agents={"strategy_choice_agent": strategy_choice_agent},
        flows={
            "strategy_choice_flow": Flow(agents=["strategy_choice_agent"], context_pass=[True])
        },
        mode=ExecutionMode.SEQUENTIAL
    )
    
    strategy = await orchestrator.process_request(input_data, flow_name="strategy_choice_flow")
    strategy = strategy['content']['strategy_choice_agent'].content
    if "The choosed strategy was: " in strategy:
        return strategy.partition('The choosed strategy was: ')[2]
    else:
        return strategy

async def learning_framework(expanded_context, strategy_choice, provider):
    df = pd.read_csv('references_data.csv')
    
    if strategy_choice in list(frameworks.keys()):
        framework_url = frameworks[strategy_choice]   
        framework_content = df[df['reference'] == framework_url]['content'].iloc[0].replace('\n','')
    else:
        framework_content = strategy_choice

    solution_generation_task = CodeSolutionTask(
        description="Generate a code to solve a scientific problem based on a framework requested.",
        goal="Given a complete description of a scientific problem and on a description of the framework desired to solve it, generate a python code to solve the problem.",
        output_format="A robust and well designed code for reproductibility.",
        llm_provider=provider
    )
    
    solver_agent = Agent(
        role=AgentRole.SPECIALIST,
        tasks=[solution_generation_task],
        llm_provider=provider
    )
    
    orchestrator = Orchestrator(
        agents={"solver_agent": solver_agent},
        flows={
            "solver_flow": Flow(agents=["solver_agent"], context_pass=[True])
        },
        mode=ExecutionMode.SEQUENTIAL
    )
    
    input_data = f"""
    Problem Description: \n{expanded_context}

    Framework to Solve it: \n{framework_content}
    """

    solution = await orchestrator.process_request(input_data, flow_name="solver_flow")
    return solution['content']['solver_agent'].content.partition('python\n')[2].partition('```')[0]

def github_search(term):
    term = term.replace('github ','').replace(' ', '%20')
    url=f'https://github.com/search?q={term}&type=repositories'
    search = requests.get(url).content
    
    contents = []
    soup = BeautifulSoup(search, 'html.parser')
    repos_container = soup.find_all('div',class_='Box-sc-g0xbh4-0 kYLlPM')
    for repo in repos_container:
        url = "https://www.github.com"+repo.find('a', class_='prc-Link-Link-85e08')['href']
        contents.append(url)
    return contents

def papers_search(term):
    term = term.replace(' ', '+')
    url = f"https://arxiv.org/search/?query={term}&searchtype=all&source=header"
    
    request = requests.get(url).content

    contents = {"url": [], "paper_name": [], "abstract": []}
    soup = BeautifulSoup(request, 'html.parser')
    papers_container = soup.find_all('li', class_='arxiv-result')
    
    for paper in papers_container:
        try:
            paper_url = paper.find('p', class_='list-title is-inline-block').find('a')['href'].replace('abs','pdf')
            name = paper.find(class_='title is-5 mathjax').text.replace('\n','').replace('  ','')
            abstract = paper.find(class_='abstract-full has-text-grey-dark mathjax').text.split('   ')[2].replace('  ','')[:-1]

            contents["url"].append(paper_url)
            contents["paper_name"].append(name)
            contents["abstract"].append(abstract)
            
        except:
            pass
    return contents

def extract_references(context, provider):   
    prompt = """
    I want to search on {source} for {format_of_data} related to this specific problem:
    \n
    {context}
    \n
    I would also like to search on {source} for {format_of_data} that use physics-informed neural networks (PINNs) to solve this specific problem.
    \n
    Additionally, I would like to search on {source} for repositories that use machine learning or deep learning in general to solve the same specific problem.
    \n
    Please return me only the terms, separated each one for a |
    """
    terms = [{
            "source": "github",
            "format_of_data": "repositories"},{
            "source": "arxiv",
            "format_of_data": "papers"}]
    
    search_terms = {"github":[],"arxiv":[]}

    for term in terms:
        prompts = provider.generate(prompt.format(
            source=term["source"],
            format_of_data=term["format_of_data"], 
            context=context))
        all_prompts = prompts.replace('\n','').split(' | ')
        search_terms[term["source"]].append(all_prompts)

    github_results, arxiv_results = [], []
    repositories_info = []
    papers_info = {"url":[],"name":[],"abstract":[]}
    
    for term in search_terms["github"][0]:
        github_results.append(github_search(term))

    for term in search_terms["arxiv"][0]:
        arxiv_results.append(papers_search(term))

    for repo in github_results:
        for url in repo:
            repositories_info.append(url)

    for paper in arxiv_results:
        urls = paper['url']
        papers_name = paper['paper_name']
        abstracts = paper['abstract']
        for i in range(len(urls)):
            url = urls[i]
            paper_name = papers_name[i]
            abstract = abstracts[i]
            papers_info["url"].append(url)
            papers_info["name"].append(paper_name)
            papers_info["abstract"].append(abstract)
    
    return repositories_info, papers_info

def extract_file_content(api_key, file_path):
    file_format = file_path.partition('.')[2]
    prompt = "Summarize the following data from a {file_format} file: \nData: {file_content}."
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-1.5-flash')
                
    if file_format == 'jpg' or file_format == 'png':
        img = Image.open(file_path)
        response = model.generate_content(img)
        return response.candidates[0].content.parts[0].text
    if file_format == 'csv':
        df = pd.read_csv(file_path)
        file_content = df.to_string()            
    if file_format == 'xls' or file_format == 'xlsx':
        df = pd.read_excel(file_path)
        file_content = df.to_string()
    if file_format == 'json':
        with open(file_path, 'r') as file:
            data = json.load(file)
        file_content = json.dumps(data, indent=4)
    if file_format == 'tiff':
        with tifffile.TiffFile(file_path) as tif:
            metadata = tif.pages[0].tags
            file_content = str(metadata)
    if file_format == 'pdf':
        file_content=''
        pdf = PyPDF2.PdfReader(file_path)
        for page in pdf.pages:
            file_content+=page.extract_text()
    if file_format == 'ppt':
        file_content = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            file_content += run.text
    if file_format == 'mp4':
        cap = cv2.VideoCapture(file_path)
        frames = []
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            frames.append(frame)
        cap.release()
        descriptions = []
        for frame in frames:
            _, buffer = cv2.imencode('.jpg', frame)
            img = buffer.tobytes()
            response = model.generate_content(img)
            descriptions.append(response.text)
        file_content = "\n".join(descriptions)
    if file_format == 'mat':
        data = scipy.io.loadmat(file_path)
        file_content = str(data)
    if file_format == 'npy' or file_format == 'npz':
        array = np.load(file_path)    
        file_content = np.array2string(array)
    if file_format == 'dxf':
        doc = ezdxf.readfile(file_path)
        msp = doc.modelspace()
        data = []
        for entity in msp:
            data.append(str(entity.dxf))
        file_content = "\n".join(data)
    if file_format == 'hdf5':
        file_content = ""
        with h5py.File(file_path, 'r') as f:
            def print_attrs(name, obj):
                nonlocal file_content
                file_content += f"Group or Dataset: {name}\n"
                for key, val in obj.attrs.items():
                    file_content += f"  Attribute: {key} = {val}\n"
                    if isinstance(obj, h5py.Dataset):
                        file_content += f"  Data: {obj[:]}\n"
            f.visititems(print_attrs)
    if file_format == 'nc':
        file_content = ""
        with Dataset(file_path, 'r') as nc:
            for var_name in nc.variables:
                var = nc.variables[var_name]
                file_content += f"Variable: {var_name}\n"
                file_content += f"  Data: {var[:]}\n"
                for attr_name in var.ncattrs():
                    file_content += f"  Attribute: {attr_name} = {var.getncattr(attr_name)}\n"
    if file_format == 'fits':
        file_content = ""
        with fits.open(file_path) as hdul:
            for hdu in hdul:
                file_content += f"Extension: {hdu.name}\n"
                file_content += f"  Header: {hdu.header}\n"
                if hdu.data is not None:
                    file_content += f"  Data: {hdu.data}\n"
    if file_format == 'parquet':
        df = pd.read_parquet(file_path)
        file_content = df.to_string()
    if file_format == 'avro':
        file_content = ""
        with open(file_path, 'rb') as fo:
            reader = avro.datafile.Reader(fo, avro.schema.ReadersWriterSchema())
            for record in reader:
                file_content += str(record) + "\n"
    if file_format == 'ply':
        file_content=''
        mesh = trimesh.load(file_path)
        file_content = f"Vertices: {mesh.vertices}\n"
        file_content += f"Faces: {mesh.faces}\n"
    if file_format == 'pcd':
        file_content=''
        pcd = o3d.io.read_point_cloud(file_path)
        points = pcd.points
        colors = pcd.colors
        file_content = f"Points: {points}\n"
        file_content += f"Colors: {colors}\n"
    prompt = prompt.format(file_format=file_format, file_content=file_content)
    return model.generate_content(prompt)
            
async def main(
        context,
        geometry = '',
        boundary_conditions = '',
        initial_conditions = '',
        equations = '',
        desired_framework = 'None',
        file_path='',
        ):
    
    test_keys = [
        "AIzaSyCoC6voLEtOEOg5caWaqEIXBh8CiYWoUaY",
        "AIzaSyDA3r3LpI8cIGm4AVoaDQ65mDMD10GNTVM"
    ]

    gemini_config = GeminiConfig(
        api_keys=test_keys,
        temperature=0.7,
        verbose=False,
        max_tokens=10000
    )
    
    provider = GeminiProvider(gemini_config)
    
    organized_context=f'Context: \n{context}'
    
    df = pd.read_csv('references_data.csv')
    
    if desired_framework in list(frameworks.keys()):
        framework_url = frameworks[desired_framework]   
        framework_content = df[df['reference'] == framework_url]['content'].iloc[0].replace('\n','')
    
    if geometry != '':
        organized_context+=f'\nGeometry: \n{geometry}'
    if boundary_conditions != '':
        organized_context+=f'\nBoundary Condition: \n{boundary_conditions}'
    if initial_conditions != '':
        organized_context+=f'\nInitial Conditions: \n{initial_conditions}'
    if equations != '':
        organized_context+=f'\nEquations: \n{equations}\n'
    if desired_frameworks != 'None':
        organized_context+=f'\nPreferential Framework: \n{framework_content}'
    else:
        framework_content = desired_framework
    if file_path != '':
        for api_key in test_keys:
            try:
                media_content = extract_file_content(api_key, file_path)
                organized_context+=f"\nMedia File Information to analyze: {media_content}"
            except:
                continue
    expanded_context = await context_expansion(organized_context, provider)
    strategy_choice = await choice_strategy(expanded_context, provider)
    
    if 'A different framework: ' in strategy_choice:
        strategy_choice = strategy_choice.replace('A different framework: ','')
    
    solution = await learning_framework(expanded_context, strategy_choice, provider)
    
    repositories, papers = extract_references(expanded_context, provider)

    output = f"**<span style='color: blue; font-size: 1.2em;'>Hi, welcome to SciML helper.</span>**\n\n"
    output += "<span style='color: green;'>Let's start trying to understand better your problem...</span>\n\n"
    output += f"<span style='color: black;'>This is the problem you want to solve:</span> \n\n<span style='color: purple;'>{organized_context.replace(framework_content, desired_framework)}</span>\n\n"
    output += f"<span style='color: black;'>My Suggestion for solving the problem using this strategy is following those code steps:</span>\n\n<span style='color: orange;'>{solution}</span>\n\n"
    output += "<span style='color: black;'>\n\nFinally, here are some references for research\n:</span>"

    output+="<span style='color: black;'>\n\nGithub Repositories:\n\n</span>"
    cont=1
    for repository in repositories:
        output+=f'{cont}. {repository}\n'
        cont+=1
    output+="<span style='color: black;'>\nOpen Papers:\n\n</span>"                  
    cont=1
    for i in range(len(papers['url'])):
        paper_name = papers['name'][i]
        url = papers['url'][i]
        abstract = papers['abstract'][i]
        output+=f'{cont}. Name: {paper_name}\nURL: {url}\nAbstract: {abstract}\n'
        cont+=1
        if cont == 11:
            break
    
    return output

def run_async_function(context, geometry, boundary_conditions, initial_conditions, equations, desired_frameworks, file_path):
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    return loop.run_until_complete(main(context, geometry, boundary_conditions, initial_conditions, equations, desired_frameworks, file_path))

st.title("SciML Helper App")
st.markdown("### Fill in the details of your problem below:")

context = st.text_area("Context (Required)", "", help="Describe the problem you want to solve.")
uploaded_file = st.file_uploader("Upload a file to increase the context", help="Upload a file to be processed.")
geometry = st.text_area("Geometry", "", help="Provide details about the problem's geometry.")
boundary_conditions = st.text_area("Boundary Conditions", "", help="Specify the boundary conditions.")
initial_conditions = st.text_area("Initial Conditions", "", help="Provide initial conditions for the problem.")
equations = st.text_area("Equations", "", help="Define the equations governing the system.")

corrected_frameworks_list = ["None"]
corrected_frameworks_list.extend(frameworks_list)
desired_frameworks = st.selectbox("Preferential Frameworks", corrected_frameworks_list, help="Select one or more frameworks you prefer.")

if st.button("Solve Problem"):
    if not context.strip():
        st.error("Context is required to proceed.")
    else:
        file_path = ''
        if uploaded_file is not None:
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file.write(uploaded_file.getvalue())
                file_path = temp_file.name
        with st.spinner("Now take a glass of Wine and Relax...\n\nThe results could take from 2 to 5 minutes."):
            result = run_async_function(context, geometry, boundary_conditions, initial_conditions, equations, desired_frameworks, file_path)
            st.markdown(result, unsafe_allow_html=True)
        if file_path:
            os.remove(file_path)