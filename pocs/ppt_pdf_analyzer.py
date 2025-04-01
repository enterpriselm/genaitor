
import streamlit as st
import os
import asyncio
import PyPDF2
import pptx
import sys

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from src.core import Orchestrator, Flow, ExecutionMode
from presets.agents import document_agent, question_agent, search_agent, response_agent

def extract_text_from_pdf(pdf_file):
    """Extracts text from a PDF file."""
    try:
        reader = PyPDF2.PdfReader(pdf_file)
        text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        return text
    except Exception as e:
        return str(e)

def extract_text_from_ppt(ppt_file):
    """Extracts text from a PPT file."""
    try:
        presentation = pptx.Presentation(ppt_file)
        text = "\n".join([
            shape.text for slide in presentation.slides
            for shape in slide.shapes if hasattr(shape, "text")
        ])
        return text
    except Exception as e:
        return str(e)

async def process_document(document_content, question):
    orchestrator = Orchestrator(
        agents={
            "document_agent": document_agent,
            "question_agent": question_agent,
            "search_agent": search_agent,
            "response_agent": response_agent
        },
        flows={
            "document_qa_flow": Flow(
                agents=["document_agent", "question_agent", "search_agent", "response_agent"],
                context_pass=[True, True, True, True]
            )
        },
        mode=ExecutionMode.SEQUENTIAL
    )
    try:
        result = await orchestrator.process_request(
            {"document_content": document_content, "user_question": question},
            flow_name='document_qa_flow'
        )
        return result
    except Exception as e:
        return {"success": False, "error": str(e)}

st.title("📄 Document Q&A System")

uploaded_file = st.file_uploader("Upload a PDF or PPTX document", type=["pdf", "pptx"])
question = st.text_input("Enter your question about the document")

if uploaded_file and question:
    st.write("Processing...")
    
    if uploaded_file.name.endswith(".pdf"):
        document_content = extract_text_from_pdf(uploaded_file)
    elif uploaded_file.name.endswith(".pptx"):
        document_content = extract_text_from_ppt(uploaded_file)
    else:
        st.error("Unsupported document format.")
        st.stop()
    
    if document_content:
        result = asyncio.run(process_document(document_content, question))
        
        if result["success"]:
            st.subheader("Answer:")
            st.write(result['content']['response_agent'].content)
        else:
            st.error(f"Error: {result['error']}")